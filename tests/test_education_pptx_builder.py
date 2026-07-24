import json
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
SKILL = ROOT / "skills" / "education-industry-observation"
LAYOUT = SKILL / "assets" / "colleague-template" / "layout.json"
TEMPLATE = SKILL / "assets" / "教育行业观察_双周报模板.pptx"
BUILDER = SKILL / "scripts" / "build_biweekly_pptx.mjs"
WRAPPER = SKILL / "scripts" / "build_biweekly_pptx.ps1"


def slide_text(slide):
    return "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))


class EducationPptxTemplateTests(unittest.TestCase):
    def test_layout_manifest_matches_colleague_visual_and_dynamic_pagination(self):
        layout = json.loads(LAYOUT.read_text(encoding="utf-8"))

        self.assertEqual("colleague-biweekly-v1", layout["template_id"])
        self.assertEqual([720, 540], layout["reference"]["page_size_pts"])
        self.assertEqual(8, layout["reference"]["page_count"])
        self.assertEqual([960, 720], layout["canvas_px"])
        self.assertEqual(15, layout["pagination"]["max_pages"])
        self.assertEqual(3, layout["pagination"]["max_items_per_page"])
        self.assertEqual(11, len(layout["reference_slots"]))
        self.assertEqual("#C3D4A7", layout["colors"]["header_green"])
        self.assertEqual("#B5CA92", layout["colors"]["label_green"])
        self.assertEqual("KaiTi", layout["fonts"]["body"]["typeface"])
        self.assertEqual(16, layout["fonts"]["body"]["font_size_pt"])

    def test_template_asset_is_eight_page_four_by_three_reference_deck(self):
        self.assertTrue(TEMPLATE.exists())
        presentation = Presentation(TEMPLATE)

        self.assertAlmostEqual(10.0, presentation.slide_width / 914400, places=2)
        self.assertAlmostEqual(7.5, presentation.slide_height / 914400, places=2)
        self.assertEqual(8, len(presentation.slides))

    def test_reference_template_uses_original_section_page_map(self):
        presentation = Presentation(TEMPLATE)
        expected = (
            "1 行业速览-上市公司",
            "1 行业速览-上市公司",
            "2 行业速览-其他",
            "2 行业速览-其他",
            "2 行业速览-其他",
            "3 行业速览-政策",
            "3 行业速览-政策",
        )
        actual = tuple(
            slide_text(slide)
            for index, slide in enumerate(presentation.slides)
            if index > 0
        )

        for text, title in zip(actual, expected):
            self.assertIn(title, text)

    def test_template_has_no_legacy_public_pages_or_labels(self):
        presentation = Presentation(TEMPLATE)
        all_text = "\n".join(slide_text(slide) for slide in presentation.slides)

        for forbidden in (
            "本期主要动态",
            "行业判断",
            "受益：",
            "风险：",
            "本期行业小结",
            "后续跟踪",
        ):
            self.assertNotIn(forbidden, all_text)

    def test_template_typography_and_shapes_match_fixed_contract(self):
        presentation = Presentation(TEMPLATE)
        cover_title = next(
            shape
            for shape in presentation.slides[0].shapes
            if hasattr(shape, "text") and shape.text == "教育行业观察"
        )
        header = next(
            shape
            for shape in presentation.slides[1].shapes
            if hasattr(shape, "text") and shape.text.startswith("1 行业速览-")
        )
        label = next(
            shape
            for shape in presentation.slides[1].shapes
            if hasattr(shape, "text") and "上\n市\n公\n司\n一" in shape.text
        )
        body = next(
            shape
            for shape in presentation.slides[1].shapes
            if hasattr(shape, "text") and "本期事项要点" in shape.text
        )

        self.assertEqual("SimSun", cover_title.text_frame.paragraphs[0].font.name)
        self.assertAlmostEqual(33.33, cover_title.text_frame.paragraphs[0].font.size.pt, places=1)
        self.assertEqual("KaiTi", header.text_frame.paragraphs[0].font.name)
        self.assertEqual(36.0, header.text_frame.paragraphs[0].font.size.pt)
        self.assertEqual("KaiTi", label.text_frame.paragraphs[0].font.name)
        self.assertEqual(18.0, label.text_frame.paragraphs[0].font.size.pt)
        self.assertEqual("KaiTi", body.text_frame.paragraphs[0].font.name)
        self.assertEqual(16.0, body.text_frame.paragraphs[0].font.size.pt)

        for slide in presentation.slides:
            for shape in slide.shapes:
                self.assertGreaterEqual(shape.left, 0)
                self.assertGreaterEqual(shape.top, 0)
                self.assertLessEqual(shape.left + shape.width, presentation.slide_width)
                self.assertLessEqual(shape.top + shape.height, presentation.slide_height)

    def test_builder_uses_artifact_tool_and_not_python_pptx(self):
        text = BUILDER.read_text(encoding="utf-8")
        wrapper = WRAPPER.read_text(encoding="utf-8")

        self.assertIn("@oai/artifact-tool", text)
        self.assertIn("PresentationFile.exportPptx", text)
        self.assertIn("report.page_plan", text)
        self.assertIn("layout.pagination.max_pages", text)
        self.assertNotIn("python-pptx", text.lower())
        self.assertNotIn("from pptx", text)
        self.assertIn("setup_artifact_tool_workspace.mjs", wrapper)


if __name__ == "__main__":
    unittest.main()
