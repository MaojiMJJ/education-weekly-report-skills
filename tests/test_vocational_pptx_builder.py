import importlib.util
import json
import unittest
from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
FIXTURES = ROOT / "tests" / "fixtures"
BUILDER_MODULE = (
    ROOT
    / "skills"
    / "vocational-education-weekly-report"
    / "scripts"
    / "build_weekly_pptx.py"
)
OUTPUT_DIR = ROOT / "tests" / "output"


def load_builder_module():
    spec = importlib.util.spec_from_file_location("vocational_pptx_builder", BUILDER_MODULE)
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def load_fixture(name):
    return json.loads((FIXTURES / name).read_text(encoding="utf-8"))


def slide_text(slide):
    return "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))


class VocationalPptxBuilderTests(unittest.TestCase):
    def setUp(self):
        self.output = OUTPUT_DIR / "vocational_valid.pptx"
        self.sources = OUTPUT_DIR / "vocational_sources.md"
        self.quality = OUTPUT_DIR / "vocational_quality.md"
        self.rejected_output = OUTPUT_DIR / "vocational_should_not_exist.pptx"
        for path in (self.output, self.sources, self.quality, self.rejected_output):
            if path.exists():
                path.unlink()

    def test_builds_four_by_three_research_deck_and_sources(self):
        builder = load_builder_module()
        report = load_fixture("vocational_valid.json")
        report["sections"][0]["items"][1]["evidence_table"] = {
            "columns": ["职业本科专业", "首批招生"],
            "rows": [["智能制造", "120人"], ["现代物流", "100人"], ["数字商务", "80人"]],
        }

        builder.build(report, self.output, self.sources, self.quality)

        self.assertTrue(self.output.exists())
        self.assertTrue(self.sources.exists())
        self.assertTrue(self.quality.exists())
        presentation = Presentation(self.output)
        self.assertAlmostEqual(10.0, presentation.slide_width / 914400, places=2)
        self.assertAlmostEqual(7.5, presentation.slide_height / 914400, places=2)
        self.assertEqual(8, len(presentation.slides))
        all_text = "\n".join(slide_text(slide) for slide in presentation.slides)
        self.assertIn("本期主要动态", slide_text(presentation.slides[1]))
        self.assertIn("行业判断", all_text)
        self.assertNotIn("后续跟踪", all_text)
        self.assertIn("主体背景", all_text)
        self.assertIn("受益", all_text)
        self.assertIn("风险", all_text)
        self.assertIn("测试省教育厅", all_text)
        self.assertIn("2026-06-16", all_text)
        self.assertIn("本期行业小结", slide_text(presentation.slides[-1]))
        self.assertTrue(
            any(
                getattr(shape, "has_table", False)
                for slide in presentation.slides
                for shape in slide.shapes
            )
        )
        table_slide = next(
            slide for slide in presentation.slides if any(getattr(shape, "has_table", False) for shape in slide.shapes)
        )
        table_shape = next(shape for shape in table_slide.shapes if getattr(shape, "has_table", False))
        analysis_label = next(
            shape for shape in table_slide.shapes if hasattr(shape, "text") and shape.text.strip() == "行业判断"
        )
        self.assertLessEqual(table_shape.top + table_shape.height, analysis_label.top)
        self.assertIn(
            "# 职教行业周报来源清单",
            self.sources.read_text(encoding="utf-8"),
        )
        self.assertIn("总分：34/40", self.quality.read_text(encoding="utf-8"))

    def test_rejects_old_news_digest_before_building(self):
        builder = load_builder_module()
        report = load_fixture("vocational_0709_failed.json")

        with self.assertRaises(Exception) as caught:
            builder.build(report, self.rejected_output)

        self.assertIn("core_insights", str(caught.exception))
        self.assertFalse(self.rejected_output.exists())

    def test_requires_sources_and_quality_delivery_paths(self):
        builder = load_builder_module()
        report = load_fixture("vocational_valid.json")

        with self.assertRaises(ValueError) as caught:
            builder.build(report, self.output)

        self.assertIn("来源清单", str(caught.exception))
        self.assertIn("质量报告", str(caught.exception))
        self.assertFalse(self.output.exists())

    def test_facts_use_compact_font_to_avoid_punctuation_widows(self):
        builder = load_builder_module()
        report = load_fixture("vocational_valid.json")

        builder.build(report, self.output, self.sources, self.quality)

        presentation = Presentation(self.output)
        first_fact = report["sections"][0]["items"][0]["facts"][0]
        facts_box = next(
            shape
            for slide in presentation.slides
            for shape in slide.shapes
            if hasattr(shape, "text") and first_fact[:-1] in shape.text
        )
        self.assertEqual(10.0, facts_box.text_frame.paragraphs[0].font.size.pt)

    def test_public_deck_uses_update_digest_and_hides_internal_tracking(self):
        builder = load_builder_module()
        report = load_fixture("vocational_valid.json")

        builder.build(report, self.output, self.sources, self.quality)

        presentation = Presentation(self.output)
        digest_text = slide_text(presentation.slides[1])
        all_text = "\n".join(slide_text(slide) for slide in presentation.slides)
        self.assertIn("本期主要动态", digest_text)
        self.assertIn("融资与交易", digest_text)
        self.assertIn("政策", digest_text)
        self.assertIn(report["sections"][0]["items"][0]["headline"], digest_text)
        self.assertIn("本期行业小结", slide_text(presentation.slides[-1]))
        for internal_label in ("本期核心观点", "证据事项", "后续跟踪", "下一期验证"):
            self.assertNotIn(internal_label, all_text)
        self.assertIn("tracking", report["sections"][0]["items"][0])


if __name__ == "__main__":
    unittest.main()
