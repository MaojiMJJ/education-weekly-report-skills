#!/usr/bin/env python3
"""根据通过质量校验的结构化 JSON 生成职业教育研究周报 PPTX。"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from report_quality import (
    SCORE_FIELDS,
    build_quality_markdown,
    build_sources_markdown,
    iter_events,
    validate_report,
)


FONT = "Microsoft YaHei"
TITLE_FONT = "Microsoft YaHei"
INK = RGBColor(31, 38, 46)
MUTED = RGBColor(91, 104, 112)
GREEN = RGBColor(117, 143, 76)
PALE_GREEN = RGBColor(218, 230, 193)
VERY_PALE_GREEN = RGBColor(241, 246, 232)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(159, 181, 119)
WARM = RGBColor(245, 240, 225)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=18,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.06,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    paragraph.font.name = FONT
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_bullets(slide, x, y, w, h, bullets, font_size=13.5, gap=5):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {bullet}"
        paragraph.space_after = Pt(gap)
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = INK
    return box


def add_header(slide, title, number=None):
    add_rect(slide, 0.42, 0.14, 9.16, 0.77, fill=PALE_GREEN, line=LINE)
    label = f"{number}  {title}" if number is not None else title
    add_textbox(slide, 0.55, 0.23, 8.75, 0.53, label, 24, True, INK, valign=MSO_ANCHOR.MIDDLE)


def source_footer(item):
    parts = []
    for source in item.get("sources") or []:
        label = source.get("name", "")
        if source.get("published_at"):
            label += f"（{source['published_at']}）"
        if label and label not in parts:
            parts.append(label)
    return "；".join(parts)


def add_cover(prs: Presentation, report: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_rect(slide, 0.55, 0.55, 8.9, 1.02, fill=PALE_GREEN, line=LINE)
    add_textbox(slide, 0.82, 0.73, 8.35, 0.62, report["title"], 30, True, INK, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, 0.8, 2.65, 8.4, 0.7, report["period"], 23, False, GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 4.72, 8.4, 0.45, "近两周职业教育重要资本事件与行业动态", 15, False, MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.75, 5.58, 4.5, 0.06, fill=GREEN, line=GREEN)


def public_section_name(name):
    label = str(name)
    for prefix in ("行业速览-", "职业教育-"):
        if label.startswith(prefix):
            return label[len(prefix):]
    return label


def digest_groups(report, limit=8):
    groups = []
    for section_index, section in enumerate(report.get("sections") or []):
        items = [item for item in section.get("items") or [] if item.get("content_role") == "event"]
        if items:
            groups.append(
                {
                    "name": public_section_name(section.get("name", "行业动态")),
                    "section_index": section_index,
                    "items": items,
                }
            )
    if sum(len(group["items"]) for group in groups) <= limit:
        return groups

    selected = set()
    first_items = []
    for group in groups:
        first = group["items"][0]
        first_items.append((sum(first["scores"].values()), group["section_index"], first["id"]))
    for _score, _section_index, item_id in sorted(first_items, key=lambda row: (-row[0], row[1]))[:limit]:
        selected.add(item_id)

    candidates = []
    for group in groups:
        for item_index, item in enumerate(group["items"]):
            if item["id"] not in selected:
                candidates.append(
                    (sum(item["scores"].values()), group["section_index"], item_index, item["id"])
                )
    for _score, _section_index, _item_index, item_id in sorted(
        candidates, key=lambda row: (-row[0], row[1], row[2])
    )[: max(0, limit - len(selected))]:
        selected.add(item_id)

    return [
        {**group, "items": [item for item in group["items"] if item["id"] in selected]}
        for group in groups
        if any(item["id"] in selected for item in group["items"])
    ]


def add_updates_slide(prs: Presentation, report: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "本期主要动态")
    y = 1.16
    for group in digest_groups(report):
        add_textbox(slide, 0.72, y, 2.4, 0.34, f"【{group['name']}】", 12.5, True, INK)
        bullet_height = max(0.5, len(group["items"]) * 0.5)
        add_bullets(
            slide,
            0.92,
            y + 0.36,
            8.25,
            bullet_height,
            [item["headline"] for item in group["items"]],
            11.0,
            3,
        )
        y += 0.36 + bullet_height + 0.18


def add_analysis_box(slide, item, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=VERY_PALE_GREEN, line=LINE)
    add_textbox(slide, x + 0.12, y + 0.08, 1.05, 0.28, "行业判断", 11, True, GREEN)
    add_textbox(slide, x + 1.05, y + 0.05, w - 1.18, h - 0.42, item["analysis"], 10.5, False, INK, valign=MSO_ANCHOR.MIDDLE)
    impact = "受益：{}  ｜  风险：{}".format(
        "、".join(item["beneficiaries"]),
        "、".join(item["risks"]),
    )
    add_textbox(slide, x + 1.05, y + h - 0.35, w - 1.18, 0.28, impact, 8.8, False, MUTED, valign=MSO_ANCHOR.MIDDLE)


def add_background_box(slide, item, x, y, w, h):
    add_textbox(slide, x, y, 0.9, h, "主体背景", 9.5, True, GREEN, valign=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + 0.86, y, w - 0.86, h, item["background"], 9.5, False, MUTED, valign=MSO_ANCHOR.MIDDLE)


def add_evidence_table(slide, evidence, x, y, w, h):
    columns = evidence.get("columns") or []
    rows = evidence.get("rows") or []
    if not columns or not rows:
        return
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    row_height = Inches(h / (len(rows) + 1))
    for row in table.rows:
        row.height = row_height
    column_width = Inches(w / len(columns))
    for column in table.columns:
        column.width = column_width
    for col_index, label in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = str(label)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PALE_GREEN
    for row_index, row in enumerate(rows, 1):
        for col_index in range(len(columns)):
            cell = table.cell(row_index, col_index)
            cell.text = str(row[col_index]) if col_index < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = FONT
                paragraph.font.size = Pt(9.5)
                paragraph.font.color.rgb = INK
                if row_index == 0:
                    paragraph.font.bold = True


def add_item_slide(prs: Presentation, section_name: str, item: dict[str, Any], number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, section_name, number)

    add_rect(slide, 0.45, 1.02, 1.08, 5.92, fill=PALE_GREEN, line=LINE)
    add_textbox(
        slide,
        0.55,
        1.19,
        0.88,
        5.55,
        item["short_title"],
        15,
        True,
        INK,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.02,
    )
    add_rect(slide, 1.62, 1.02, 7.96, 5.92, fill=WHITE, line=LINE)

    metadata = f"主体：{item['subject']}  ｜  日期：{item['event_date']}  ｜  类型：{item['event_type']}"
    add_textbox(slide, 1.85, 1.12, 7.45, 0.28, metadata, 9.5, True, GREEN, valign=MSO_ANCHOR.MIDDLE)

    evidence = item.get("evidence_table")
    if evidence:
        add_bullets(slide, 1.85, 1.48, 7.48, 0.96, item["facts"][:3], 10.0, 2)
        add_background_box(slide, item, 1.9, 2.5, 7.35, 0.42)
        add_evidence_table(slide, evidence, 2.0, 3.0, 7.15, 1.2)
        add_analysis_box(slide, item, 1.86, 4.38, 7.45, 1.98)
    else:
        add_bullets(slide, 1.85, 1.5, 7.46, 2.12, item["facts"], 10.0, 5)
        add_background_box(slide, item, 1.9, 3.73, 7.35, 0.48)
        add_analysis_box(slide, item, 1.86, 4.38, 7.45, 1.98)

    source = source_footer(item)
    add_textbox(slide, 0.5, 7.05, 7.5, 0.24, f"来源：{source}", 8.5, False, MUTED, valign=MSO_ANCHOR.MIDDLE)
    score = sum(item["scores"][field] for field in SCORE_FIELDS)
    add_textbox(slide, 8.15, 7.05, 1.28, 0.24, f"价值评分 {score}/25", 8.5, True, GREEN, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs: Presentation, report: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "本期行业小结")
    add_rect(slide, 0.72, 1.35, 8.56, 3.35, fill=VERY_PALE_GREEN, line=LINE)
    add_textbox(slide, 1.0, 1.7, 8.0, 2.65, report["weekly_judgment"], 17, False, INK, valign=MSO_ANCHOR.MIDDLE)
    event_count = sum(1 for _section_name, _item in iter_events(report))
    add_textbox(slide, 0.8, 5.28, 8.4, 0.45, f"本期共纳入 {event_count} 项经核验行业动态", 16, True, GREEN, align=PP_ALIGN.CENTER)
    categories = "  ｜  ".join(group["name"] for group in digest_groups(report))
    add_textbox(slide, 0.8, 5.86, 8.4, 0.42, categories, 13, False, MUTED, align=PP_ALIGN.CENTER)


def build(
    report: dict[str, Any],
    output: Path,
    sources_output: Path | None = None,
    quality_output: Path | None = None,
) -> None:
    validate_report(report)
    if sources_output is None or quality_output is None:
        raise ValueError("正式交付必须同时提供来源清单和质量报告输出路径")
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)

    add_cover(presentation, report)
    add_updates_slide(presentation, report)
    number = 1
    for section_name, item in iter_events(report):
        add_item_slide(presentation, section_name, item, number)
        number += 1
    add_summary_slide(presentation, report)

    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)

    sources_output = Path(sources_output)
    sources_output.parent.mkdir(parents=True, exist_ok=True)
    sources_output.write_text(build_sources_markdown(report), encoding="utf-8")

    quality_output = Path(quality_output)
    quality_output.parent.mkdir(parents=True, exist_ok=True)
    quality_output.write_text(build_quality_markdown(report), encoding="utf-8")


def self_test_spec() -> dict[str, Any]:
    items = []
    for index in range(1, 6):
        items.append(
            {
                "id": f"T{index}",
                "content_role": "event",
                "headline": f"[自检样例] 职业教育行业事项 {index}",
                "short_title": f"自检事项 {index}",
                "event_date": f"2026-07-0{index}",
                "subject": f"自检主体 {index}",
                "event_type": "产业学院" if index % 2 else "政策",
                "period_trigger": {
                    "type": "material_business_update" if index % 2 else "policy_issued",
                    "description": (
                        f"自检主体 {index} 于2026年7月{index}日签约落地产教项目"
                        if index % 2
                        else f"自检主体 {index} 于2026年7月{index}日发布本期政策"
                    ),
                    "source_url": f"https://www.moe.gov.cn/self-test/vocational/{index}",
                },
                "background": "该主体已具备职业教育办学、企业合作或政策执行基础，本事项用于验证研究字段和版式容量。",
                "facts": [
                    "该事项披露了明确的主体、时间和新增动作。",
                    "公开材料提供了产品、客户或政策执行信息。",
                    "事项可以形成可核验的业务或政策指标。",
                ],
                "analysis": "该事项用于验证生成器能否把职教事实与行业解释分层呈现，而不是输出新闻摘录。",
                "beneficiaries": ["职业院校和产业企业"],
                "risks": ["后续落地不及预期"],
                "tracking": "跟踪下一期公开披露的项目进度和量化指标。",
                "scores": {
                    "industry_impact": 4,
                    "policy_importance": 3,
                    "commercial_value": 4,
                    "technology_relevance": 4,
                    "investment_relevance": 3,
                },
                "score_reasons": {
                    "industry_impact": "影响明确教育场景。",
                    "policy_importance": "与政策方向相关。",
                    "commercial_value": "存在可验证业务价值。",
                    "technology_relevance": "包含技术应用。",
                    "investment_relevance": "可继续跟踪量化指标。",
                },
                "sources": [
                    {
                        "name": f"自检来源 {index}",
                        "url": f"https://www.moe.gov.cn/self-test/vocational/{index}",
                        "published_at": f"2026-07-0{index}",
                        "source_type": "test",
                        "is_primary": True,
                        "access_status": "verified",
                        "access_checked_at": "2026-07-05",
                    }
                ],
            }
        )
    return {
        "title": "职教行业周报",
        "period": "2026.06.22-2026.07.05",
        "core_insights": [
            {
                "claim": "职教自检样例验证研究结构能够稳定映射到幻灯片",
                "evidence_ids": ["T1", "T2"],
            }
        ],
        "sections": [
            {"name": "职业教育-产教融合", "items": items[:2]},
            {"name": "职业教育-政策", "items": items[2:4]},
            {"name": "职业教育-院校", "items": items[4:]},
        ],
        "weekly_judgment": "本期样例覆盖职业教育政策、院校和企业进展，事项均包含明确时间、主体、公开事实和行业影响。报告按栏目同步正式进展，并通过来源、评分和报告期门槛排除检索过程、空结果与低价值信息，使外发内容保持简洁、完整和可核验。",
        "quality_review": {
            "information_quality": {"score": 8, "reason": "职教自检事项均包含日期、主体、来源和事实字段。"},
            "analysis_depth": {"score": 8, "reason": "职教事项包含判断、受益方和风险分析。"},
            "readability": {"score": 8, "reason": "自检页面使用固定结构并限制文本容量。"},
            "strategic_value": {"score": 8, "reason": "自检内容可以形成下一期量化验证任务。"},
        },
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="根据研究型 JSON 生成职教行业周报 PPTX")
    parser.add_argument("input_json", nargs="?", help="报告 JSON 文件路径")
    parser.add_argument("--output", required=True, help="输出 PPTX 路径")
    parser.add_argument("--sources-output", required=True, help="输出 Markdown 来源清单路径")
    parser.add_argument("--quality-output", required=True, help="输出 Markdown 质量报告路径")
    parser.add_argument("--self-test", action="store_true", help="使用内置自检数据")
    args = parser.parse_args()

    if args.self_test:
        report = self_test_spec()
    else:
        if not args.input_json:
            raise SystemExit("未使用 --self-test 时必须提供 input_json")
        report = json.loads(Path(args.input_json).read_text(encoding="utf-8"))

    build(
        report,
        Path(args.output),
        Path(args.sources_output),
        Path(args.quality_output),
    )
    print(f"已写入 {args.output}")
    print(f"已写入 {args.sources_output}")
    print(f"已写入 {args.quality_output}")


if __name__ == "__main__":
    main()
