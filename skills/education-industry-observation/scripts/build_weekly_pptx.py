#!/usr/bin/env python3
"""根据通过质量校验的结构化 JSON 生成教育行业研究周报 PPTX。"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.insert(0, str(SCRIPT_DIR))

from report_quality import (
    SCORE_FIELDS,
    build_quality_markdown,
    build_sources_markdown,
    iter_briefs,
    iter_events,
    validate_report,
)


FONT = "KaiTi"
TITLE_FONT = "KaiTi"
TITLE_SIZE = 36
SIDE_LABEL_SIZE = 18
BODY_SIZE = 16
INK = RGBColor(31, 38, 46)
MUTED = RGBColor(91, 104, 112)
GREEN = RGBColor(117, 143, 76)
PALE_GREEN = RGBColor(218, 230, 193)
VERY_PALE_GREEN = RGBColor(241, 246, 232)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(159, 181, 119)
WARM = RGBColor(245, 240, 225)


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text="",
    font_size=BODY_SIZE,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.06,
    font_name=FONT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = str(text)
    paragraph.alignment = align
    paragraph.font.name = font_name
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill=WHITE, line=LINE, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.8)
    return shape


def add_bullets(slide, x, y, w, h, bullets, font_size=BODY_SIZE, gap=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for index, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {bullet}"
        paragraph.alignment = PP_ALIGN.JUSTIFY
        paragraph.space_after = Pt(gap)
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = INK
    return box


def add_paragraphs(
    slide,
    x,
    y,
    w,
    h,
    paragraphs,
    font_size=BODY_SIZE,
    gap=7,
    color=INK,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    for index, value in enumerate(paragraphs):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(value)
        paragraph.alignment = PP_ALIGN.JUSTIFY
        paragraph.space_after = Pt(gap if index < len(paragraphs) - 1 else 0)
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = color
    return box


def add_header(slide, title, number=None):
    add_rect(slide, 0.42, 0.12, 9.16, 0.82, fill=PALE_GREEN, line=LINE)
    label = f"{number}  {title}" if number is not None else title
    add_textbox(
        slide,
        0.55,
        0.16,
        8.75,
        0.7,
        label,
        TITLE_SIZE,
        False,
        INK,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.01,
        font_name=TITLE_FONT,
    )


def source_footer(item):
    parts = []
    for source in item.get("sources") or []:
        label = source.get("name", "")
        if source.get("published_at"):
            label += f"（{source['published_at']}）"
        if label and label not in parts:
            parts.append(label)
    return "；".join(parts)


def add_cover(prs: Presentation, report: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_rect(slide, 0.55, 0.55, 8.9, 1.02, fill=PALE_GREEN, line=LINE)
    add_textbox(
        slide,
        0.82,
        0.67,
        8.35,
        0.74,
        report["title"],
        TITLE_SIZE,
        False,
        INK,
        valign=MSO_ANCHOR.MIDDLE,
        font_name=TITLE_FONT,
    )
    add_textbox(slide, 0.8, 2.65, 8.4, 0.7, report["period"], 23, False, GREEN, align=PP_ALIGN.CENTER)
    add_textbox(slide, 0.8, 4.72, 8.4, 0.45, "近两周教育行业重要资本事件与行业动态", 15, False, MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 2.75, 5.58, 4.5, 0.06, fill=GREEN, line=GREEN)


def public_section_name(name):
    label = str(name)
    for prefix in ("行业速览-", "职业教育-"):
        if label.startswith(prefix):
            return label[len(prefix):]
    return label


def digest_groups(report, limit=12):
    groups = []
    for section_index, section in enumerate(report.get("sections") or []):
        items = [
            item
            for item in section.get("items") or []
            if item.get("content_role") in {"event", "brief"}
        ]
        if items:
            groups.append(
                {
                    "name": public_section_name(section.get("name", "行业动态")),
                    "section_index": section_index,
                    "items": items,
                }
            )
    if sum(len(group["items"]) for group in groups) <= limit:
        return groups

    selected = set()
    first_items = []
    for group in groups:
        first = group["items"][0]
        first_items.append(
            (
                sum((first.get("scores") or {}).values()),
                group["section_index"],
                first["id"],
            )
        )
    for _score, _section_index, item_id in sorted(first_items, key=lambda row: (-row[0], row[1]))[:limit]:
        selected.add(item_id)

    candidates = []
    for group in groups:
        for item_index, item in enumerate(group["items"]):
            if item["id"] not in selected:
                candidates.append(
                    (
                        sum((item.get("scores") or {}).values()),
                        group["section_index"],
                        item_index,
                        item["id"],
                    )
                )
    for _score, _section_index, _item_index, item_id in sorted(
        candidates, key=lambda row: (-row[0], row[1], row[2])
    )[: max(0, limit - len(selected))]:
        selected.add(item_id)

    return [
        {**group, "items": [item for item in group["items"] if item["id"] in selected]}
        for group in groups
        if any(item["id"] in selected for item in group["items"])
    ]


def add_updates_slide(prs: Presentation, report: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "本期主要动态")
    y = 1.16
    for group in digest_groups(report):
        add_textbox(slide, 0.72, y, 2.4, 0.38, f"【{group['name']}】", BODY_SIZE, False, INK)
        bullet_height = max(0.7, len(group["items"]) * 0.82)
        add_bullets(
            slide,
            0.92,
            y + 0.4,
            8.25,
            bullet_height,
            [item["headline"] for item in group["items"]],
            BODY_SIZE,
            7,
        )
        y += 0.4 + bullet_height + 0.12


def vertical_label_text(value):
    tokens = re.findall(r"[A-Za-z0-9]+|[^\s]", str(value))
    return "\n".join(tokens)


def end_sentence(value):
    return str(value).rstrip("。；;") + "。"


def add_brief_slides(prs: Presentation, report: dict[str, Any]):
    grouped = []
    for section in report.get("sections") or []:
        briefs = [
            item
            for item in section.get("items") or []
            if item.get("content_role") == "brief"
        ]
        if briefs:
            grouped.append((public_section_name(section.get("name", "行业动态")), briefs))

    for section_name, briefs in grouped:
        for start in range(0, len(briefs), 2):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = WHITE
            add_header(slide, f"行业速览 - {section_name}")
            for index, item in enumerate(briefs[start : start + 2]):
                y = 1.2 + index * 2.88
                add_rect(slide, 0.72, y, 8.56, 2.62, fill=VERY_PALE_GREEN, line=LINE, radius=True)
                add_textbox(
                    slide,
                    0.94,
                    y + 0.13,
                    7.95,
                    0.34,
                    f"{item['short_title']}  ｜  {item['event_date']}  ｜  {item['subject']}",
                    15,
                    True,
                    GREEN,
                )
                add_bullets(
                    slide,
                    1.0,
                    y + 0.52,
                    7.9,
                    1.08,
                    item["facts"],
                    BODY_SIZE,
                    3,
                )
                add_textbox(
                    slide,
                    1.0,
                    y + 1.65,
                    7.9,
                    0.52,
                    f"为什么重要：{item['why_it_matters']}",
                    BODY_SIZE,
                    False,
                    INK,
                    align=PP_ALIGN.JUSTIFY,
                )
                add_textbox(
                    slide,
                    1.0,
                    y + 2.24,
                    7.9,
                    0.2,
                    f"来源：{source_footer(item)}",
                    8.5,
                    False,
                    MUTED,
                )


def add_analysis_box(slide, item, x, y, w, h):
    add_rect(slide, x, y, w, h, fill=VERY_PALE_GREEN, line=LINE)
    add_paragraphs(
        slide,
        x + 0.14,
        y + 0.08,
        w - 0.28,
        h - 0.16,
        [
            f"行业判断：{item['analysis']}",
            f"受益：{'、'.join(item['beneficiaries'])}",
            f"风险：{'、'.join(item['risks'])}",
        ],
        BODY_SIZE,
        4,
        INK,
        MSO_ANCHOR.MIDDLE,
        margin=0.03,
    )


def add_background_box(slide, item, x, y, w, h):
    add_textbox(
        slide,
        x,
        y,
        w,
        h,
        f"主体背景：{item['background']}",
        BODY_SIZE,
        False,
        INK,
        align=PP_ALIGN.JUSTIFY,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.03,
    )


def add_evidence_table(slide, evidence, x, y, w, h):
    columns = evidence.get("columns") or []
    rows = evidence.get("rows") or []
    if not columns or not rows:
        return
    table_shape = slide.shapes.add_table(len(rows) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    row_height = Inches(h / (len(rows) + 1))
    for row in table.rows:
        row.height = row_height
    column_width = Inches(w / len(columns))
    for column in table.columns:
        column.width = column_width
    for col_index, label in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = str(label)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PALE_GREEN
    for row_index, row in enumerate(rows, 1):
        for col_index in range(len(columns)):
            cell = table.cell(row_index, col_index)
            cell.text = str(row[col_index]) if col_index < len(row) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
    for row_index, row in enumerate(table.rows):
        for cell in row.cells:
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.name = FONT
                paragraph.font.size = Pt(BODY_SIZE)
                paragraph.font.color.rgb = INK
                if row_index == 0:
                    paragraph.font.bold = False


def add_financing_content(slide, item):
    details = item["financing_details"]
    funding = [
        f"投融资：{details['round']}；{details['amount']}（{details['currency']}）；"
        f"投资方：{'、'.join(details['investors'])}。",
        f"资金用途：{'、'.join(details['fund_use'])}。",
        f"业务定位：{details['business_positioning']}；领域：{details['business_domain']}；"
        f"{details['customer_side']}；{details['delivery_mode']}；"
        f"{details['offering_type']}；用户类型：{details['user_type']}。",
    ]
    company = [
        f"公司概况：成立于{details['founded_at']}；用户{details['users']}；门店{details['stores']}；"
        f"直营/加盟：{details['store_model']}；此前投资方：{'、'.join(details['prior_investors'])}。",
        f"财务情况：流水{details['gross_billing']}；营收{details['revenue']}；利润{details['profit']}。",
    ]
    add_paragraphs(slide, 1.85, 1.48, 7.45, 2.08, funding, BODY_SIZE, 6)
    add_paragraphs(slide, 1.85, 3.64, 7.45, 1.16, company, BODY_SIZE, 5)
    add_analysis_box(slide, item, 1.82, 4.9, 7.5, 1.72)


def add_cooperation_content(slide, item):
    details = item["cooperation_details"]
    cooperation = [
        f"合作主体：{'、'.join(details['parties'])}；主体类型：{'、'.join(details['party_types'])}。",
        f"新增动作：{details['action_type']}；内容：{details['cooperation_content']}。",
        f"业务/学校：{details['business_or_school']}；地区：{details['location']}。",
    ]
    execution = [
        f"实施计划：{details['implementation_plan']}。",
        f"商业化路径：{details['commercialization']}。",
    ]
    add_paragraphs(slide, 1.85, 1.48, 7.45, 2.1, cooperation, BODY_SIZE, 6)
    add_paragraphs(slide, 1.85, 3.68, 7.45, 1.12, execution, BODY_SIZE, 6)
    add_analysis_box(slide, item, 1.82, 4.9, 7.5, 1.72)


def add_policy_content(slide, item):
    details = item["policy_details"]
    source = [
        end_sentence(
            f"政策来源：{details['issuing_body']}；《{details['policy_name']}》；"
            f"发布于{details['issued_at']}；生效时间：{details['effective_at']}"
        ),
        end_sentence(f"执行范围：{details['scope_level']}，{details['scope_description']}"),
    ]
    clause_parts = []
    if details["prohibited_rules"]:
        clause_parts.append(end_sentence(f"禁止性规定：{'；'.join(details['prohibited_rules'])}"))
    if details["restrictive_requirements"]:
        clause_parts.append(end_sentence(f"限制性要求：{'；'.join(details['restrictive_requirements'])}"))
    if details["supportive_measures"]:
        clause_parts.append(end_sentence(f"倡导性/支持性内容：{'；'.join(details['supportive_measures'])}"))
    add_paragraphs(slide, 1.85, 1.48, 7.45, 1.24, source, BODY_SIZE, 6)
    add_paragraphs(slide, 1.85, 2.8, 7.45, 1.74, clause_parts, BODY_SIZE, 6)
    add_analysis_box(slide, item, 1.82, 4.64, 7.5, 1.98)


def add_generic_content(slide, item):
    evidence = item.get("evidence_table")
    if evidence:
        add_bullets(slide, 1.85, 1.48, 7.48, 1.5, item["facts"][:3], BODY_SIZE, 5)
        add_evidence_table(slide, evidence, 1.95, 3.08, 7.25, 1.22)
        add_analysis_box(slide, item, 1.82, 4.42, 7.5, 2.2)
    else:
        add_bullets(slide, 1.85, 1.48, 7.48, 1.96, item["facts"], BODY_SIZE, 6)
        add_background_box(slide, item, 1.85, 3.54, 7.45, 0.9)
        add_analysis_box(slide, item, 1.82, 4.54, 7.5, 2.08)


def add_item_slide(prs: Presentation, section_name: str, item: dict[str, Any], number: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    display_section = str(section_name).replace("行业速览-", "行业速览 - ")
    add_header(slide, display_section, number)

    add_rect(slide, 0.45, 1.02, 1.08, 5.92, fill=PALE_GREEN, line=LINE)
    add_textbox(
        slide,
        0.55,
        1.19,
        0.88,
        5.55,
        vertical_label_text(item["short_title"]),
        SIDE_LABEL_SIZE,
        False,
        INK,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        margin=0.02,
    )
    add_rect(slide, 1.62, 1.02, 7.96, 5.92, fill=WHITE, line=LINE)

    metadata = f"主体：{item['subject']}  ｜  日期：{item['event_date']}  ｜  类型：{item['event_type']}"
    add_textbox(slide, 1.85, 1.12, 7.45, 0.28, metadata, 9.5, True, GREEN, valign=MSO_ANCHOR.MIDDLE)

    if item.get("financing_details"):
        add_financing_content(slide, item)
    elif item.get("cooperation_details"):
        add_cooperation_content(slide, item)
    elif item.get("policy_details"):
        add_policy_content(slide, item)
    else:
        add_generic_content(slide, item)

    source = source_footer(item)
    add_textbox(slide, 0.5, 7.05, 7.5, 0.24, f"来源：{source}", 8.5, False, MUTED, valign=MSO_ANCHOR.MIDDLE)
    score = sum(item["scores"][field] for field in SCORE_FIELDS)
    add_textbox(slide, 8.15, 7.05, 1.28, 0.24, f"价值评分 {score}/25", 8.5, True, GREEN, align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)


def add_summary_slide(prs: Presentation, report: dict[str, Any]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_header(slide, "本期行业小结")
    add_rect(slide, 0.72, 1.35, 8.56, 3.35, fill=VERY_PALE_GREEN, line=LINE)
    add_textbox(
        slide,
        1.0,
        1.7,
        8.0,
        2.65,
        report["weekly_judgment"],
        BODY_SIZE,
        False,
        INK,
        align=PP_ALIGN.JUSTIFY,
        valign=MSO_ANCHOR.MIDDLE,
    )
    event_count = sum(1 for _section_name, _item in iter_events(report))
    brief_count = sum(1 for _section_name, _item in iter_briefs(report))
    add_textbox(
        slide,
        0.8,
        5.28,
        8.4,
        0.45,
        f"本期共纳入 {event_count + brief_count} 项经核验行业动态，其中 {event_count} 项深度分析",
        16,
        True,
        GREEN,
        align=PP_ALIGN.CENTER,
    )
    categories = "  ｜  ".join(group["name"] for group in digest_groups(report))
    add_textbox(slide, 0.8, 5.86, 8.4, 0.42, categories, 13, False, MUTED, align=PP_ALIGN.CENTER)


def build(
    report: dict[str, Any],
    output: Path,
    sources_output: Path | None = None,
    quality_output: Path | None = None,
) -> None:
    validate_report(report)
    if sources_output is None or quality_output is None:
        raise ValueError("正式交付必须同时提供来源清单和质量报告输出路径")
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(7.5)

    add_cover(presentation, report)
    add_updates_slide(presentation, report)
    add_brief_slides(presentation, report)
    number = 1
    for section_name, item in iter_events(report):
        add_item_slide(presentation, section_name, item, number)
        number += 1
    add_summary_slide(presentation, report)

    output = Path(output)
    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output)

    sources_output = Path(sources_output)
    sources_output.parent.mkdir(parents=True, exist_ok=True)
    sources_output.write_text(build_sources_markdown(report), encoding="utf-8")

    quality_output = Path(quality_output)
    quality_output.parent.mkdir(parents=True, exist_ok=True)
    quality_output.write_text(build_quality_markdown(report), encoding="utf-8")


def self_test_spec() -> dict[str, Any]:
    items = []
    for index in range(1, 6):
        items.append(
            {
                "id": f"T{index}",
                "content_role": "event",
                "headline": f"[自检样例] 教育行业事项 {index}",
                "short_title": f"自检事项 {index}",
                "event_date": f"2026-07-0{index}",
                "subject": f"自检主体 {index}",
                "event_type": "产品" if index % 2 else "政策",
                "period_trigger": {
                    "type": "product_launched" if index % 2 else "policy_issued",
                    "description": f"自检主体 {index} 于2026年7月{index}日发布本期新增事项",
                    "source_url": f"https://www.moe.gov.cn/self-test/{index}",
                },
                "background": "该主体已具备教育产品、客户或政策执行基础，本事项用于验证研究字段和版式容量。",
                "policy_details": (
                    {
                        "issuing_body": f"自检主管部门 {index}",
                        "policy_name": f"自检教育政策 {index}",
                        "issued_at": f"2026-07-0{index}",
                        "effective_at": f"2026-07-0{index}",
                        "scope_level": "地方性",
                        "scope_description": "适用于自检地区教育机构和学校。",
                        "prohibited_rules": ["禁止违反规定开展教育服务"],
                        "restrictive_requirements": ["按要求完成备案和信息披露"],
                        "supportive_measures": ["支持符合条件的教育数字化应用"],
                    }
                    if index % 2 == 0
                    else None
                ),
                "facts": [
                    "该事项披露了明确的主体、时间和新增动作。",
                    "公开材料提供了产品、客户或政策执行信息。",
                    "事项可以形成可核验的业务或政策指标。",
                ],
                "analysis": "该事项用于验证生成器能否把事实与行业解释分层呈现，而不是输出新闻摘录。",
                "beneficiaries": ["教育服务企业"],
                "risks": ["后续落地不及预期"],
                "tracking": "跟踪下一期公开披露的项目进度和量化指标。",
                "scores": {
                    "industry_impact": 4,
                    "policy_importance": 3,
                    "commercial_value": 4,
                    "technology_relevance": 4,
                    "investment_relevance": 3,
                },
                "score_reasons": {
                    "industry_impact": "影响明确教育场景。",
                    "policy_importance": "与政策方向相关。",
                    "commercial_value": "存在可验证业务价值。",
                    "technology_relevance": "包含技术应用。",
                    "investment_relevance": "可继续跟踪量化指标。",
                },
                "sources": [
                    {
                        "name": f"自检来源 {index}",
                        "url": f"https://www.moe.gov.cn/self-test/{index}",
                        "published_at": f"2026-07-0{index}",
                        "source_type": "test",
                        "is_primary": True,
                        "access_status": "verified",
                        "access_checked_at": "2026-07-05",
                    }
                ],
            }
        )
    return {
        "title": "教育行业观察",
        "period": "2026.06.22-2026.07.05",
        "core_insights": [
            {
                "claim": "自检样例验证研究结构能够稳定映射到幻灯片",
                "evidence_ids": ["T1", "T2"],
            }
        ],
        "sections": [
            {"name": "行业速览-产品", "items": items[:2]},
            {"name": "行业速览-政策", "items": items[2:4]},
            {"name": "行业速览-AI教育", "items": items[4:]},
        ],
        "weekly_judgment": "本期样例覆盖政策与产品发布，事项均包含明确时间、主体、公开事实和行业影响。报告按栏目同步正式进展，并通过来源、评分和报告期门槛排除检索过程、空结果与低价值信息。各事项的受益对象和主要风险均有公开依据，外发内容聚焦已经发生的变化，使读者可以快速了解行业进展及其影响。",
        "quality_review": {
            "information_quality": {"score": 8, "reason": "自检事项均包含日期、主体、来源和事实字段。"},
            "analysis_depth": {"score": 8, "reason": "自检事项包含判断、受益方和风险分析。"},
            "readability": {"score": 8, "reason": "自检页面使用固定结构并限制文本容量。"},
            "strategic_value": {"score": 8, "reason": "自检内容能够清楚呈现本期事项及其行业影响。"},
        },
    }


def main() -> None:
    parser = argparse.ArgumentParser(description="根据研究型 JSON 生成教育行业观察 PPTX")
    parser.add_argument("input_json", nargs="?", help="报告 JSON 文件路径")
    parser.add_argument("--output", required=True, help="输出 PPTX 路径")
    parser.add_argument("--sources-output", required=True, help="输出 Markdown 来源清单路径")
    parser.add_argument("--quality-output", required=True, help="输出 Markdown 质量报告路径")
    parser.add_argument("--self-test", action="store_true", help="使用内置自检数据")
    args = parser.parse_args()

    if args.self_test:
        report = self_test_spec()
    else:
        if not args.input_json:
            raise SystemExit("未使用 --self-test 时必须提供 input_json")
        report = json.loads(Path(args.input_json).read_text(encoding="utf-8"))

    build(
        report,
        Path(args.output),
        Path(args.sources_output),
        Path(args.quality_output),
    )
    print(f"已写入 {args.output}")
    print(f"已写入 {args.sources_output}")
    print(f"已写入 {args.quality_output}")


if __name__ == "__main__":
    main()
